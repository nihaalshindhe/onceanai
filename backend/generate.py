from fastapi import APIRouter, HTTPException, Header, Body, Depends
from fastapi.responses import StreamingResponse
from sqlalchemy.orm import Session
from database import get_db
from models import Project, Section, Revision
from schemas import ProjectCreate, Refine
from llm import generate_text
import jwt, os, io
from docx import Document
from pptx import Presentation

router = APIRouter()
SECRET = os.getenv("JWT_SECRET")

def get_user_id(token: str):
    try:
        data = jwt.decode(token, SECRET, algorithms=["HS256"])
        return data["id"]
    except:
        raise HTTPException(status_code=401, detail="Invalid token")

@router.get("/sections/{project_id}")
def get_sections(project_id: int, authorization: str = Header(...), db: Session = Depends(get_db)):
    token = authorization.replace("Bearer ", "")
    user_id = get_user_id(token)
    sections = db.query(Section).join(Project).filter(
        Section.project_id == project_id,
        Project.user_id == user_id
    ).all()
    return sections

@router.get("/projects")
def get_projects(authorization: str = Header(...), db: Session = Depends(get_db)):
    token = authorization.replace("Bearer ", "")
    user_id = get_user_id(token)
    projects = db.query(Project).filter(Project.user_id == user_id).all()
    return projects

@router.post("/generate-doc")
def generate_doc(data: ProjectCreate = Body(...), authorization: str = Header(...), db: Session = Depends(get_db)):
    token = authorization.replace("Bearer ", "")
    user_id = get_user_id(token)

    project = Project(
        user_id=user_id,
        topic=data.topic,
        doc_type=data.doc_type,
        outline="\n".join([s.title for s in data.sections])
    )
    db.add(project)
    db.commit()
    db.refresh(project)

    sections_content = []
    for s in data.sections:
        final_prompt = f"{s.prompt}\n\nRemove all comments. Add the best option. Return only the asked content."
        content = generate_text(final_prompt)
        section = Section(
            project_id=project.id,
            title=s.title,
            content=content
        )
        db.add(section)
        sections_content.append((s.title, content))
    db.commit()

    if data.doc_type == "docx":
        doc = Document()
        doc.add_heading(data.topic, level=0)
        for title, content in sections_content:
            doc.add_heading(title, level=1)
            doc.add_paragraph(content)
        tmp = io.BytesIO()
        doc.save(tmp)
        tmp.seek(0)
        return StreamingResponse(
            tmp,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": f"attachment; filename={data.topic}.docx"}
        )

    elif data.doc_type == "pptx":
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = data.topic
        slide.placeholders[1].text = "Overview"

        slide_layout = prs.slide_layouts[1]
        for title, content in sections_content:
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            slide.placeholders[1].text = content

        tmp = io.BytesIO()
        prs.save(tmp)
        tmp.seek(0)
        return StreamingResponse(
            tmp,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename={data.topic}.pptx"}
        )

    else:
        raise HTTPException(status_code=400, detail="Unsupported document type")

@router.post("/refine")
def refine_section(data: Refine = Body(...), authorization: str = Header(...), db: Session = Depends(get_db)):
    token = authorization.replace("Bearer ", "")
    user_id = get_user_id(token)

    section = db.query(Section).join(Project).filter(
        Section.id == data.section_id,
        Project.user_id == user_id
    ).first()
    if not section:
        raise HTTPException(status_code=404, detail="Section not found")

    refine_prompt = (
        f"Refine the following section **without adding any extra headings, comments, or markers**.\n\n"
        f"Original Content:\n{section.content}\n\n"
        f"Instruction:\n{data.prompt}\n\n"
        f"Return only the refined content."
    )

    new_text = generate_text(refine_prompt)

    revision = Revision(
        section_id=data.section_id,
        prompt=data.prompt,
        result=new_text,
        feedback=data.feedback,
        comment=data.comment
    )
    db.add(revision)
    db.commit()

    return {"updated_content": new_text}

@router.post("/refine/approve")
def approve_refinement(section_id: int = Body(...), doc_type: str = Body(...), authorization: str = Header(...), db: Session = Depends(get_db)):
    token = authorization.replace("Bearer ", "")
    user_id = get_user_id(token)

    section = db.query(Section).join(Project).filter(
        Section.id == section_id,
        Project.user_id == user_id
    ).first()
    if not section:
        raise HTTPException(status_code=404, detail="Section not found")

    latest_revision = db.query(Revision).filter(
        Revision.section_id == section_id
    ).order_by(Revision.id.desc()).first()
    if not latest_revision:
        raise HTTPException(status_code=400, detail="No refinement found")

    section.content = latest_revision.result
    db.commit()

    project = db.query(Project).filter(Project.id == section.project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    sections = db.query(Section).filter(Section.project_id == project.id).order_by(Section.id).all()

    if doc_type == "docx":
        doc = Document()
        doc.add_heading(project.topic, level=0)
        for sec in sections:
            doc.add_heading(sec.title, level=1)
            doc.add_paragraph(sec.content)
        tmp = io.BytesIO()
        doc.save(tmp)
        tmp.seek(0)
        return StreamingResponse(
            tmp,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": f"attachment; filename={project.topic}.docx"}
        )

    elif doc_type == "pptx":
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = project.topic

        slide_layout = prs.slide_layouts[1]
        for sec in sections:
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = sec.title
            slide.placeholders[1].text = sec.content

        tmp = io.BytesIO()
        prs.save(tmp)
        tmp.seek(0)
        return StreamingResponse(
            tmp,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename={project.topic}.pptx"}
        )

    else:
        raise HTTPException(status_code=400, detail="Unsupported document type")
