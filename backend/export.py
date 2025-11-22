from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
from database import SessionLocal
from models import Project, Section
from docx import Document
from pptx import Presentation

router = APIRouter()

@router.get("/docx/{project_id}")
def export_docx(project_id: int):
    db = SessionLocal()

    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    sections = db.query(Section).filter(Section.project_id == project_id).all()

    doc = Document()
    doc.add_heading(project.topic, level=0)

    for s in sections:
        doc.add_heading(s.title, level=1)
        doc.add_paragraph(s.content)

    filename = f"project_{project_id}.docx"
    doc.save(filename)

    return FileResponse(filename, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", filename=filename)

@router.get("/pptx/{project_id}")
def export_pptx(project_id: int):
    db = SessionLocal()

    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    sections = db.query(Section).filter(Section.project_id == project_id).all()

    pres = Presentation()

    for s in sections:
        slide = pres.slides.add_slide(pres.slide_layouts[1])
        slide.shapes.title.text = s.title
        slide.placeholders[1].text = s.content

    filename = f"project_{project_id}.pptx"
    pres.save(filename)

    return FileResponse(filename, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename=filename)
